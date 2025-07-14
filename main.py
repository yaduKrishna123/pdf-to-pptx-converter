from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader
import os
import uuid
import threading

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

tasks = {}

def pdf_to_ppt(task_id, pdf_path, pptx_path='output.pptx', dpi=200):
    reader = PdfReader(pdf_path)
    total_pages = len(reader.pages)

    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    for i in range(1, total_pages + 1):
        images = convert_from_path(pdf_path, dpi=dpi, first_page=i, last_page=i)
        img = images[0]

        temp_image = f'temp_{uuid.uuid4().hex}.png'
        img.save(temp_image, 'PNG')

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(temp_image, 0, 0, width=prs.slide_width, height=prs.slide_height)

        os.remove(temp_image)

        # ✅ Update progress
        tasks[task_id]['status'] = f"converted {i}/{total_pages}"

    prs.save(pptx_path)

def convert_in_background(task_id, input_path, output_path):
    try:
        pdf_to_ppt(task_id, input_path, output_path)
        tasks[task_id]['status'] = 'done'
        tasks[task_id]['output_path'] = output_path
    except Exception as e:
        tasks[task_id] = {
            'status': 'error',
            'message': str(e)
        }
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

@app.after_request
def add_cors_headers(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type,Authorization'
    response.headers['Access-Control-Allow-Methods'] = 'GET,POST,OPTIONS'
    return response

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    task_id = str(uuid.uuid4())
    input_pdf_path = f"{task_id}.pdf"
    output_pptx_path = f"{task_id}.pptx"

    file.save(input_pdf_path)
    tasks[task_id] = {'status': 'processing'}

    thread = threading.Thread(target=convert_in_background, args=(task_id, input_pdf_path, output_pptx_path))
    thread.start()

    return jsonify({'task_id': task_id}), 202

@app.route('/status/<task_id>', methods=['GET'])
def status(task_id):
    task = tasks.get(task_id)
    if not task:
        return jsonify({'error': 'Task not found'}), 404
    return jsonify(task)

@app.route('/download/<task_id>', methods=['GET'])
def download(task_id):
    task = tasks.get(task_id)
    if not task or task['status'] != 'done':
        return jsonify({'error': 'File not ready or task not found'}), 404

    output_path = task['output_path']
    if os.path.exists(output_path):
        response = send_file(output_path, as_attachment=True)
        os.remove(output_path)
        del tasks[task_id]
        return response
    return jsonify({'error': 'Output file missing'}), 404

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
