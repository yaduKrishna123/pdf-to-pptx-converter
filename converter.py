from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from flask import Flask, request, send_file, jsonify
import os

def pdf_to_ppt(pdf_path, pptx_path='output.pptx',dpi=200):
    images = convert_from_path(pdf_path,dpi=dpi)

    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    for img in images:
        temp_image = 'temp_slide.png'
        img.save(temp_image,'PNG')

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.add_picture(temp_image, 0, 0, width=prs.slide_width,height = prs.slide_height)
        os.remove(temp_image)

    prs.save(pptx_path)
    print(f"power point saved to {pptx_path}")

pdf_to_ppt("joicy.pdf","outputjoicy.pptx")
      